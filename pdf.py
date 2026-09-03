from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper to add title and bullet slides
def add_title_slide(title_text, subtitle_text, author_text, supervisor_text, dept_text, date_text):
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 58, 95)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.7))
    tf2 = sub_box.text_frame
    tf2.text = subtitle_text
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Authors
    auth_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5))
    tf3 = auth_box.text_frame
    tf3.text = author_text
    tf3.paragraphs[0].font.size = Pt(22)
    tf3.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Supervisor
    sup_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.5))
    tf4 = sup_box.text_frame
    tf4.text = supervisor_text
    tf4.paragraphs[0].font.size = Pt(20)
    tf4.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    tf4.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Department
    dept_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.5))
    tf5 = dept_box.text_frame
    tf5.text = dept_text
    tf5.paragraphs[0].font.size = Pt(18)
    tf5.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    tf5.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.5))
    tf6 = date_box.text_frame
    tf6.text = date_text
    tf6.paragraphs[0].font.size = Pt(18)
    tf6.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    tf6.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Footer & page number
    add_footer(slide, 1)

def add_content_slide(title_text, bullet_points, slide_num):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 58, 95)

    # Add blue line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(74, 144, 226)
    line.line.fill.background()

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    for point in bullet_points:
        p = tf2.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
        p.level = 0

    add_footer(slide, slide_num)

def add_footer(slide, slide_num):
    # University text at bottom right
    footer = slide.shapes.add_textbox(Inches(10), Inches(7.1), Inches(3), Inches(0.3))
    tf = footer.text_frame
    tf.text = "Ho Technical University"
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Page number at bottom center
    page = slide.shapes.add_textbox(Inches(6), Inches(7.1), Inches(1), Inches(0.3))
    tf2 = page.text_frame
    tf2.text = f"Slide {slide_num} of 6"
    tf2.paragraphs[0].font.size = Pt(10)
    tf2.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 1: Title Slide
add_title_slide(
    "Enhanced Adaptive Hybrid Collaborative Filtering Recommendation System (AHCF)\nwith Bisecting K-Means Clustering for Multi-Vendor E-commerce Platforms",
    "Project Defense",
    "Wonder Awyiwu (0322080289) & Godwin Teye Tetteh (0322080402)",
    "Supervisor: Mr. Robert Agboyi",
    "Department of Computer Science, Ho Technical University",
    "8th – 9th May 2026"
)

# Slide 2: Introduction
intro_bullets = [
    "Recommender systems help users discover products based on behaviour and ratings",
    "Cold-start problem: new users/items have no history → poor recommendations",
    "Multi-vendor platforms in Ghana (Jumia, Konga) face constant influx of new users and products",
    "AHCF (Agboyi, 2019): already solves cold-start, but retraining is computationally expensive – 'algorithms are computationally expensive' (Agboyi, 2019, p.94)",
    "Our solution: Replace standard k-means with bisecting k-means to reduce retraining time"
]
add_content_slide("Introduction", intro_bullets, 2)

# Slide 3: Problem Statement
prob_bullets = [
    "PROBLEM: Retraining overhead in dynamic multi-vendor environments",
    "AHCF retrains 50% faster than offline models, but still slow for real-time updates",
    "In multi-vendor sites: new users daily, new products constantly, imbalanced data (some vendors have 1,000 products, some 10)",
    "Standard k-means recomputes ALL clusters from scratch → time-consuming",
    "Shared hosting in Ghana (limited CPU/RAM) makes delays worse → users abandon site",
    "",
    "RESEARCH GAP:",
    "   • Bisecting k-means is proven faster in weblog and image segmentation",
    "   • But never applied to AHCF for e-commerce",
    "   • Agboyi (2019, p.95) recommended future work: 'a model that is 100% online'"
]
add_content_slide("Problem Statement", prob_bullets, 3)

# Slide 4: Aim and Objectives
aim_bullets = [
    "Aim: Enhance AHCF with bisecting k-means to improve computational efficiency and reduce retraining time",
    "",
    "Objectives:",
    "   1. Replace standard k-means with bisecting k-means in the AHCF clustering component",
    "   2. Integrate the enhanced engine into a multi-vendor Laravel platform",
    "   3. Compare training time of standard k-means vs bisecting k-means on real data",
    "   4. Ensure recommendation quality remains high"
]
add_content_slide("Aim and Objectives", aim_bullets, 4)

# Slide 5: Methodology
method_bullets = [
    "DATA: 102 users, 207 products, 942 ratings from live Laravel site",
    "",
    "BISECTING K-MEANS (5 clusters):",
    "   • Build user-product matrix (rows=users, columns=products)",
    "   • Start with all users in one cluster → repeatedly split largest cluster into two",
    "   • Stop when 5 clusters achieved",
    "",
    "API ENDPOINTS: POST /train (training time) | GET /recommend/<user_id> (product IDs)",
    "",
    "LARAVEL INTEGRATION: RecommendationService calls API; fallback to behaviour-based recommendations if API unavailable",
    "",
    "PERFORMANCE COMPARISON:",
    "   • Standard k-means training time: ≈1.2 seconds",
    "   • Bisecting k-means training time: ≈0.5 seconds",
    "   • ≈58% reduction; recommendation quality unchanged (precision@10 ≈0.73)"
]
add_content_slide("Methodology", method_bullets, 5)

# Slide 6: Thank You
thank_bullets = [
    "Questions Welcome",
    "",
    "We appreciate your feedback and suggestions.",
    "",
    "Contact: wonderawyiwu19@gmail.com"
]
add_content_slide("Thank You", thank_bullets, 6)

# Save the file
prs.save("project_defense_slides.pptx")
print("✅ Presentation saved as 'project_defense_slides.pptx'")